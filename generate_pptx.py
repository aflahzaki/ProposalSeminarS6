from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from copy import deepcopy

# Constants
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
LIGHT_GREEN = RGBColor(0xC6, 0xEF, 0xCE)
GRAY = RGBColor(0x80, 0x80, 0x80)
LIGHT_GRAY = RGBColor(0xD9, 0xD9, 0xD9)
FONT_NAME = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title(slide, text, left=Inches(0.5), top=Inches(0.3), width=Inches(12), height=Inches(0.8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT_NAME
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    return txBox

def add_textbox(slide, text, left, top, width, height, font_size=Pt(18), bold=False, alignment=PP_ALIGN.LEFT, color=BLACK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT_NAME
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox

def add_placeholder_box(slide, text, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    shape.line.color.rgb = GRAY
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT_NAME
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER
    return shape

def add_content_paragraphs(slide, items, left, top, width, height, font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.name = FONT_NAME
        p.font.size = font_size
        p.font.color.rgb = BLACK
        p.space_after = Pt(8)
    return txBox

# ============ SLIDE 1: Cover ============
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
add_textbox(slide, "SEMINAR PROPOSAL", Inches(1), Inches(0.8), Inches(11.333), Inches(0.6),
            font_size=Pt(20), bold=True, alignment=PP_ALIGN.CENTER, color=DARK_BLUE)
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.333), Inches(2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Analisis Preskriptif Kualitas Air Menggunakan Natural Gradient Boosting dan Counterfactual Explanations"
p.font.name = FONT_NAME
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
p.alignment = PP_ALIGN.CENTER

info_items = [
    "",
    "Aflah Zaki Siregar",
    "NIM: 103062300095",
    "Program Studi: S1 IT-KJ-23-002",
    "",
    "Dosen Pembimbing 1: Nurul Ilmi, S.Kom., M.T.",
    "Dosen Pembimbing 2: Rana Zaini Fathiyana, S.ST., M.T."
]
for item in info_items:
    p = tf.add_paragraph()
    p.text = item
    p.font.name = FONT_NAME
    p.font.size = Pt(18)
    p.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 2: Latar Belakang (1/2) ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "LATAR BELAKANG")
items = [
    "1. Kualitas air minum yang memenuhi standar adalah aspek fundamental kesehatan masyarakat. 2 miliar orang belum memiliki akses air minum aman [1][2].",
    "",
    "2. Permenkes No. 2/2023 menetapkan Standar Baku Mutu Kesehatan Lingkungan (SBMKL). Monitoring laboratorium tradisional membutuhkan biaya tinggi dan waktu lama [3].",
    "",
    "3. Machine learning (XGBoost, RF) mampu memprediksi kualitas air secara cepat, NAMUN hanya menghasilkan prediksi DETERMINISTIK \u2014 label 'layak/tidak layak' tanpa informasi ketidakpastian [4][5]."
]
add_content_paragraphs(slide, items, Inches(0.5), Inches(1.2), Inches(8.5), Inches(5.5))
add_placeholder_box(slide, "[VISUALISASI: Masukkan gambar infografis masalah air / statistik WHO \u2014 buat sendiri di Canva]",
                    Inches(9.2), Inches(4.5), Inches(3.8), Inches(2.5))

# ============ SLIDE 3: Latar Belakang (2/2) ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "LATAR BELAKANG (Lanjutan)")
items = [
    "1. NGBoost memodelkan DISTRIBUSI PROBABILITAS penuh \u2192 menghasilkan P(y=1|x) terkalibrasi yang memberitahu SEBERAPA YAKIN prediksinya [6].",
    "",
    "2. Namun mengetahui 'tidak layak 89%' saja BELUM CUKUP \u2014 perlu tahu APA YANG HARUS DIPERBAIKI.",
    "",
    "3. Explainable AI (DiCE) menghasilkan rekomendasi PRESKRIPTIF: 'turunkan pH dari 9.9 ke 8.4, kurangi Chloramines ke 3.5 mg/l' [7].",
    "",
    "4. Integrasi NGBoost + DiCE = framework PRESKRIPTIF LENGKAP: prediksi + keyakinan + rekomendasi perbaikan [6][7]."
]
add_content_paragraphs(slide, items, Inches(0.5), Inches(1.2), Inches(8.5), Inches(5.5))
add_placeholder_box(slide, "[VISUALISASI: Figure 1 dari Duan et al. 2020 \u2014 file: NGBoost_Natural_Gradient_Boosting_for_Probabilistic_Prediction.pdf, Halaman 1, pojok kanan atas]",
                    Inches(9.2), Inches(1.2), Inches(3.8), Inches(2.5))
add_placeholder_box(slide, "[VISUALISASI: Box perbandingan 'XGBoost: yes/no only' vs 'NGBoost+DiCE: probability + rekomendasi' \u2014 buat sendiri di Canva]",
                    Inches(9.2), Inches(4.5), Inches(3.8), Inches(2.5))

# ============ SLIDE 4: Penelitian Terkait & Gap ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "PENELITIAN TERKAIT & GAP")

# Create table
rows, cols = 6, 5
tbl_left = Inches(0.5)
tbl_top = Inches(1.3)
tbl_width = Inches(12.3)
tbl_height = Inches(4.0)
table_shape = slide.shapes.add_table(rows, cols, tbl_left, tbl_top, tbl_width, tbl_height)
table = table_shape.table

# Set column widths
col_widths = [Inches(2.8), Inches(2.5), Inches(2.2), Inches(2.5), Inches(2.3)]
for i, w in enumerate(col_widths):
    table.columns[i].width = w

# Header
headers = ["Author & Tahun", "Metode", "Tipe Prediksi", "Explainability", "Domain"]
for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    p = cell.text_frame.paragraphs[0]
    p.font.name = FONT_NAME
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_BLUE

# Data rows
data = [
    ["Aslam et al. (2022)", "Hybrid NN+XGB", "Deterministik", "Deskriptif", "Air"],
    ["Al Bataineh et al. (2026)", "XGBoost+NN", "Deterministik", "Pasif (SHAP)", "Air"],
    ["Park et al. (2022)", "XGBoost+SHAP", "Deterministik", "Pasif (SHAP)", "Air"],
    ["Mothilal et al. (2020)", "DiCE", "Deterministik", "Preskriptif", "Non-Air"],
    ["Penelitian Ini", "NGBoost+DiCE", "Probabilistik", "Aktif Preskriptif", "Air"],
]

for row_idx, row_data in enumerate(data, start=1):
    for col_idx, val in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.font.name = FONT_NAME
        p.font.size = Pt(14)
        p.font.color.rgb = BLACK
        if row_idx == 5:  # "Penelitian Ini" row
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GREEN
            p.font.bold = True

# Gap text
add_textbox(slide, "GAP: Belum ada integrasi model PROBABILISTIK + counterfactual explanations di domain kualitas air",
            Inches(0.5), Inches(5.8), Inches(12), Inches(0.8), font_size=Pt(18), bold=True, color=DARK_BLUE)

# ============ SLIDE 5: Perumusan Masalah ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "PERUMUSAN MASALAH")

rqs = [
    "RQ1: Bagaimana performa NGBoost dalam memodelkan kelayakan air secara PROBABILISTIK dibandingkan baseline (XGBoost, RF)?",
    "RQ2: Bagaimana implementasi DiCE pada NGBoost untuk menghasilkan REKOMENDASI PRESKRIPTIF perubahan parameter air?",
    "RQ3: Sejauh mana kualitas rekomendasi counterfactual memenuhi properti validity, proximity, sparsity, diversity?"
]

for i, rq in enumerate(rqs):
    top = Inches(1.5 + i * 2.0)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.5), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFE)
    shape.line.color.rgb = DARK_BLUE
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = rq
    p.font.name = FONT_NAME
    p.font.size = Pt(18)
    p.font.color.rgb = BLACK

# ============ SLIDE 6: Tujuan Penelitian ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "TUJUAN PENELITIAN")
items = [
    "1. Menganalisis performa NGBoost secara komparatif terhadap baseline dari metrik klasifikasi (Accuracy, Precision, Recall, F1) DAN kalibrasi (NLL, ECE).",
    "",
    "2. Mengimplementasikan DiCE pada NGBoost untuk menghasilkan rekomendasi perubahan parameter fisikokimia yang mengubah status air.",
    "",
    "3. Menganalisis kualitas counterfactual berdasarkan validity, proximity, sparsity, diversity."
]
add_content_paragraphs(slide, items, Inches(0.5), Inches(1.3), Inches(12), Inches(5.5))

# ============ SLIDE 7: Batasan Masalah ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "BATASAN MASALAH")
items = [
    "\u25cf Dataset: Water Potability (Kadiwal) + Canada (Nature Scientific Data)",
    "",
    "\u25cf Parameter: 8-9 fitur fisikokimia",
    "",
    "\u25cf Klasifikasi: Binary (Kadiwal) + 5-class CCME (Canada)",
    "",
    "\u25cf Lingkungan: Python, Google Colab"
]
add_content_paragraphs(slide, items, Inches(0.5), Inches(1.3), Inches(12), Inches(4.5))
add_placeholder_box(slide, "[JADWAL KEGIATAN: Dikosongkan \u2014 isi sendiri]",
                    Inches(0.5), Inches(5.8), Inches(12), Inches(1.2))

# ============ SLIDE 8: Data Penelitian ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "DATA PENELITIAN")

# Left column
left_items = [
    "Dataset Utama: Kadiwal",
    "\u2022 3,276 sampel",
    "\u2022 9 fitur fisikokimia",
    "\u2022 Binary (Potable/Not Potable)",
    "\u2022 Missing values: pH, Sulfate, THMs",
    "\u2022 Sumber: Kaggle"
]
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5.5), Inches(4.0))
tf = txBox.text_frame
tf.word_wrap = True
for i, item in enumerate(left_items):
    if i == 0:
        p = tf.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(20)
    else:
        p = tf.add_paragraph()
        p.font.size = Pt(18)
    p.text = item
    p.font.name = FONT_NAME
    p.font.color.rgb = BLACK
    p.space_after = Pt(6)

# Right column
right_items = [
    "Dataset Pendukung: Canada",
    "\u2022 3,949 sampel",
    "\u2022 8 fitur fisikokimia",
    "\u2022 5-class (CCME WQI)",
    "\u2022 0 missing values",
    "\u2022 Sumber: Nature Scientific Data (2025)"
]
txBox = slide.shapes.add_textbox(Inches(6.8), Inches(1.3), Inches(5.5), Inches(4.0))
tf = txBox.text_frame
tf.word_wrap = True
for i, item in enumerate(right_items):
    if i == 0:
        p = tf.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(20)
    else:
        p = tf.add_paragraph()
        p.font.size = Pt(18)
    p.text = item
    p.font.name = FONT_NAME
    p.font.color.rgb = BLACK
    p.space_after = Pt(6)

add_placeholder_box(slide, "[VISUALISASI: Bar chart distribusi kelas kedua dataset \u2014 dari Colab output Cell 5 notebook Kadiwal dan Canada]",
                    Inches(0.5), Inches(5.8), Inches(12), Inches(1.4))

# ============ SLIDE 9: Alur Metodologi ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "ALUR METODOLOGI")

# Box 1
shape1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.8), Inches(3.8), Inches(3.5))
shape1.fill.solid()
shape1.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFE)
shape1.line.color.rgb = DARK_BLUE
tf = shape1.text_frame
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = "PREPROCESSING"
p.font.name = FONT_NAME
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
p.alignment = PP_ALIGN.CENTER
for item in ["\u2022 Split 70/15/15", "\u2022 Median Imputation", "\u2022 MinMax Scaling", "\u2022 SMOTE-ENN evaluated"]:
    p = tf.add_paragraph()
    p.text = item
    p.font.name = FONT_NAME
    p.font.size = Pt(14)
    p.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.CENTER

# Arrow 1
add_textbox(slide, "\u2192", Inches(4.2), Inches(3.0), Inches(0.6), Inches(0.6), font_size=Pt(36), bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

# Box 2
shape2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(1.8), Inches(3.8), Inches(3.5))
shape2.fill.solid()
shape2.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFE)
shape2.line.color.rgb = DARK_BLUE
tf = shape2.text_frame
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = "FASE A: PEMODELAN"
p.font.name = FONT_NAME
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
p.alignment = PP_ALIGN.CENTER
for item in ["\u2022 NGBoost (Bernoulli)", "\u2022 XGBoost (baseline)", "\u2022 RF (baseline)", "\u2022 Evaluasi + McNemar"]:
    p = tf.add_paragraph()
    p.text = item
    p.font.name = FONT_NAME
    p.font.size = Pt(14)
    p.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.CENTER

# Arrow 2
add_textbox(slide, "\u2192", Inches(8.7), Inches(3.0), Inches(0.6), Inches(0.6), font_size=Pt(36), bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

# Box 3
shape3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.3), Inches(1.8), Inches(3.8), Inches(3.5))
shape3.fill.solid()
shape3.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFE)
shape3.line.color.rgb = DARK_BLUE
tf = shape3.text_frame
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = "FASE B: DiCE"
p.font.name = FONT_NAME
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
p.alignment = PP_ALIGN.CENTER
for item in ["\u2022 Counterfactual generation", "\u2022 Constraint WHO/Permenkes", "\u2022 Evaluasi CF properti", "\u2022 Output preskriptif"]:
    p = tf.add_paragraph()
    p.text = item
    p.font.name = FONT_NAME
    p.font.size = Pt(14)
    p.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.CENTER

# Bottom note
add_textbox(slide, "Zero Leakage Methodology \u2014 Split SEBELUM preprocessing [ref: van de Mortel 2025]",
            Inches(0.5), Inches(6.0), Inches(12), Inches(0.6), font_size=Pt(16), bold=True, color=DARK_BLUE)

# ============ SLIDE 10: Metrik Evaluasi ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "METRIK EVALUASI & LINGKUNGAN SIMULASI")

# Table
rows, cols = 4, 3
table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(12.3), Inches(3.0))
table = table_shape.table
table.columns[0].width = Inches(2.5)
table.columns[1].width = Inches(5.5)
table.columns[2].width = Inches(4.3)

headers = ["Aspek", "Metrik", "Tujuan"]
for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    p = cell.text_frame.paragraphs[0]
    p.font.name = FONT_NAME
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_BLUE

data = [
    ["Klasifikasi", "Accuracy, Precision, Recall, F1, AUC-ROC", "Performa diskriminatif"],
    ["Kalibrasi", "NLL, ECE, Calibration Curve", "Kualitas probabilitas"],
    ["Counterfactual", "Validity, Proximity, Sparsity, Diversity", "Kualitas rekomendasi"],
]
for row_idx, row_data in enumerate(data, start=1):
    for col_idx, val in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.font.name = FONT_NAME
        p.font.size = Pt(14)
        p.font.color.rgb = BLACK

add_textbox(slide, "Lingkungan: Python 3.10+ | Google Colab | ngboost | dice-ml | xgboost | scikit-learn",
            Inches(0.5), Inches(5.5), Inches(12), Inches(0.8), font_size=Pt(18), color=BLACK)

# ============ SLIDE 11: Hasil Sementara ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "HASIL SEMENTARA")
add_placeholder_box(slide, "[DIKOSONGKAN \u2014 Isi sendiri dengan hasil dari Colab]",
                    Inches(0.5), Inches(2.0), Inches(12), Inches(3.5))
add_textbox(slide, "Saran isi: Tabel accuracy Canada (96%) vs Kadiwal (67-70%), contoh output DiCE, uncertainty zone analysis",
            Inches(0.5), Inches(6.0), Inches(12), Inches(0.8), font_size=Pt(14), color=GRAY)

# ============ SLIDE 12: Terima Kasih ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_textbox(slide, "TERIMA KASIH", Inches(1), Inches(2.5), Inches(11.333), Inches(1.5),
            font_size=Pt(44), bold=True, alignment=PP_ALIGN.CENTER, color=DARK_BLUE)
add_textbox(slide, "Aflah Zaki Siregar \u2014 103062300095", Inches(1), Inches(4.2), Inches(11.333), Inches(0.8),
            font_size=Pt(22), alignment=PP_ALIGN.CENTER, color=BLACK)
add_textbox(slide, "Pertanyaan?", Inches(1), Inches(5.5), Inches(11.333), Inches(0.8),
            font_size=Pt(20), alignment=PP_ALIGN.CENTER, color=GRAY)

# Save
output_path = "/projects/sandbox/ProposalSeminarS6/PPT_Proposal_Seminar_AflahZaki.pptx"
prs.save(output_path)
print(f"PPTX saved to: {output_path}")
